import os
import logging
import langchain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
# from langchain.document_loaders import PyPDFLoader
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.schema import Document
from docx import Document as DocxDocument
from pptx import Presentation
from io import BytesIO

# Setting up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

DB_FAISS_PATH = "vectorstores7/db_faiss/"

# Functions to process documents from bytes
def load_docx_from_bytes(file_bytes):
    logger.info("Processing DOCX bytes")
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error processing DOCX bytes: {e}")
        return ""

def load_pptx_from_bytes(file_bytes):
    logger.info("Processing PPTX bytes")
    try:
        presentation = Presentation(BytesIO(file_bytes))
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error processing PPTX bytes: {e}")
        return ""

def create_vector_db_from_memory(file_bytes, file_type):
    try:
        documents = []

        # Process document based on type
        if file_type == "pdf":
            loader = PyPDFLoader(BytesIO(file_bytes))
            documents.extend(loader.load())
        elif file_type == "docx":
            text = load_docx_from_bytes(file_bytes)
            if text.strip():
                documents.append(Document(page_content=text))
        elif file_type == "pptx":
            text = load_pptx_from_bytes(file_bytes)
            if text.strip():
                documents.append(Document(page_content=text))

        logger.info(f"Loaded {len(documents)} documents from bytes.")
        
        if not documents:
            logger.error("No documents found in the provided bytes.")
            return

        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        texts = text_splitter.split_documents(documents)
        logger.info(f"Total text chunks created: {len(texts)}")
        
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2",
                                           model_kwargs={'device': 'cpu'})
        
        db = FAISS.from_documents(texts, embeddings)
        db.save_local(DB_FAISS_PATH)
        logger.info(f"FAISS vector store saved to {DB_FAISS_PATH}")
    
    except Exception as e:
        logger.error(f"Error during vector DB creation from memory: {e}")

# Example usage: replace with your actual file bytes and type
if __name__ == "__main__":
    # Simulate file bytes and type (use actual bytes and type in your application)
    file_bytes = b"example file content"
    file_type = "pdf"  # or "docx", "pptx"
    create_vector_db_from_memory(file_bytes, file_type)
